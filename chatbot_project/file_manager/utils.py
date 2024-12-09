from pptx import Presentation
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
import faiss


nltk.download('punkt_tab')
nltk.download('stopwords')

from transformers import pipeline
from sentence_transformers import SentenceTransformer

# Load the embedding model globally
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')  # Lightweight and efficient model


from PyPDF2 import PdfReader

def extract_text_from_pdf(file_path):
    """
    Extracts text from a PDF file.
    """
    reader = PdfReader(file_path)
    text_content = []

    for page in reader.pages:
        text_content.append(page.extract_text())

    return "\n".join(text_content)


def extract_text_from_ppt(file_path):
    presentation = Presentation(file_path)
    text_content = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text_content.append(paragraph.text)

    return "\n".join(text_content)

def preprocess_text(text):
    # Tokenize text
    tokens = word_tokenize(text.lower())
    # Remove stopwords
    stop_words = set(stopwords.words('english'))
    filtered_tokens = [word for word in tokens if word.isalnum() and word not in stop_words]
    return filtered_tokens


def chunk_text(text, max_tokens=300):
    """
    Splits text into chunks within the token limit.
    """
    chunks = []
    paragraphs = text.split("\n\n")  # Split by paragraphs
    current_chunk = []
    current_length = 0
    for paragraph in paragraphs:
        paragraph_length = len(paragraph.split())
        if current_length + paragraph_length > max_tokens:
            chunks.append("\n\n".join(current_chunk))
            current_chunk = [paragraph]
            current_length = paragraph_length
        else:
            current_chunk.append(paragraph)
            current_length += paragraph_length

    if current_chunk:
        chunks.append("\n\n".join(current_chunk))

    return chunks

def generate_embeddings(chunks):
    """
    Generate embeddings for a list of text chunks.
    """
    embeddings = embedding_model.encode(chunks, show_progress_bar=True)
    return embeddings

import numpy as np

def preprocess_and_store_chunks(text):
    """
    Preprocess text by chunking, generating embeddings, and storing the chunks with embeddings.
    """
    # Step 1: Chunk the text
    chunks = chunk_text(text)

    # Step 2: Generate embeddings for the chunks
    embeddings = generate_embeddings(chunks)

    # Step 3: Store chunks and embeddings (for now, as in-memory lists)
    return chunks, np.array(embeddings)


import openai
import os

def generate_detailed_answer(context, question):
    """
    Use OpenAI's GPT-4 or GPT-3.5-turbo to generate a detailed answer based on the context and question.
    """
    openai.api_key = os.getenv("OPENAI_API_KEY")

    try:
        response = openai.ChatCompletion.create(
            model="gpt-4",  # You can also use "gpt-3.5-turbo" for lower cost
            messages=[
                {"role": "system", "content": "You are a WealthTech expert providing detailed, conversational, and sales-friendly answers."},
                {"role": "user", "content": f"Context: {context}\n\nQuestion: {question}"}
            ],
            temperature=0.7
        )
        # Extract and return the content of the response
        return response['choices'][0]['message']['content']
    except Exception as e:
        return f"Error generating response: {e}"


def create_faiss_index(embeddings):
    """
    Create a FAISS index from embeddings.
    """
    # Step 1: Define the dimension of the embeddings
    dimension = embeddings.shape[1]
    
    # Step 2: Create a FAISS index for L2 similarity (default)
    index = faiss.IndexFlatL2(dimension)
    
    # Step 3: Add embeddings to the index
    # pylint: disable=no-value-for-parameter
    index.add(embeddings)
    return index

import pickle

def save_faiss_index_and_chunks(index, chunks, file_prefix):
    """
    Save the FAISS index and the corresponding chunks to disk.
    """
    # Save the FAISS index
    faiss.write_index(index, f"{file_prefix}_faiss.index")
    
    # Save the chunks as a pickle file
    with open(f"{file_prefix}_chunks.pkl", "wb") as f:
        pickle.dump(chunks, f)

def load_faiss_index_and_chunks(file_prefix):
    """
    Load the FAISS index and corresponding chunks from disk.
    """
    # Load the FAISS index
    index = faiss.read_index(f"{file_prefix}_faiss.index")
    
    # Load the chunks
    with open(f"{file_prefix}_chunks.pkl", "rb") as f:
        chunks = pickle.load(f)
    
    return index, chunks

'''def search_index(query, index, chunks, top_k=3):
    """
    Search the FAISS index for the top_k most relevant chunks.
    """
    # Generate embedding for the query
    query_embedding = embedding_model.encode([query])
    
    # Search the FAISS index
    distances, indices = index.search(np.array(query_embedding), top_k)
    
    # Retrieve top_k chunks and their distances
    results = [(chunks[idx], distances[0][i]) for i, idx in enumerate(indices[0])]
    return results
'''
from .models import UploadedFile

def load_indices_and_chunks():
    """
    Load all FAISS indices and corresponding chunks from the database.
    """
    indices = []
    all_chunks = []
    # pylint: disable=no-member
    for ppt in UploadedFile.objects.all():
        index, chunks = load_faiss_index_and_chunks(ppt.name)
        indices.append(index)
        all_chunks.extend(chunks)

    return indices, all_chunks

def query_faiss_across_files(question, indices, chunks, top_k=3):
    """
    Query across multiple FAISS indices and return top results.
    """
    results = []
    for index, file_chunks in zip(indices, chunks):
        file_results = search_index(question, index, file_chunks, top_k)
        results.extend(file_results)

    # Sort by similarity score
    results = sorted(results, key=lambda x: x[1], reverse=True)
    return results[:top_k]

def search_index(question, index, chunks, top_k=3):
    """
    Search a single FAISS index and retrieve top results.
    """
    question_embedding = embedding_model.encode([question])
    distances, indices = index.search(question_embedding, top_k)
    results = [(chunks[idx], distances[i]) for i, idx in enumerate(indices[0]) if idx < len(chunks)]
    return results

def find_relevant_chunks(question, index, chunks, top_k=3):
    """
    Retrieve the most relevant chunks for a given question using the FAISS index.
    """
    # Generate embeddings for the question
    question_embedding = embedding_model.encode([question])

    # Perform FAISS search
    distances, indices = index.search(question_embedding, top_k)

    # Retrieve the corresponding chunks
    relevant_chunks = [chunks[idx] for idx in indices[0] if idx < len(chunks)]

    return relevant_chunks


from dotenv import load_dotenv
import os

load_dotenv()  # Load environment variables
print(f"Loaded API Key: {os.getenv('OPENAI_API_KEY')}")  # Debug
